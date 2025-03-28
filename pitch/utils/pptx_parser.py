from pptx import Presentation
from typing import List, Dict

def parse_pptx(filepath: str) -> List[Dict]:
    """Parse PPTX file into structured slides data"""
    slides = []
    
    try:
        prs = Presentation(filepath)
        
        for i, slide in enumerate(prs.slides):
            content = {
                'headings': [],
                'paragraphs': [],
                'bullet_points': [],
                'tables': []
            }
            
            title = f"Slide {i + 1}"
            
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                    
                text = shape.text.strip()
                if not text:
                    continue
                
                # Detect title shape (usually first shape or has specific formatting)
                if (i == 0 and not content['headings']) or shape.is_title:
                    title = text
                    content['headings'].append(text)
                # Detect bullet points
                elif any(p.level > 0 for p in shape.text_frame.paragraphs):
                    content['bullet_points'].append(text)
                else:
                    content['paragraphs'].append(text)
            
            slides.append({
                "slide_number": i + 1,
                "title": title,
                "content": content,
                "meta": {
                    "layout": slide.slide_layout.name,
                    "has_notes": slide.has_notes_slide,
                    "shape_count": len(slide.shapes)
                }
            })
            
    except Exception as e:
        raise ValueError(f"Failed to parse PPTX: {str(e)}")
    
    return slides